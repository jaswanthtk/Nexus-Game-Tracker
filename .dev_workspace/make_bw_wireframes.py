import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_bw_blueprint_pptx():
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors: Black & White Architectural Blueprint
    BG_COLOR = RGBColor(255, 255, 255)       # Crisp White background
    TEXT_DARK = RGBColor(20, 20, 20)         # Charcoal Black
    TEXT_MUTED = RGBColor(100, 100, 100)     # Blueprint Gray
    BOX_BG = RGBColor(248, 249, 250)         # Off-white / Light gray fill
    BORDER_COLOR = RGBColor(40, 40, 40)      # Solid structural border
    ACCENT_BOX = RGBColor(235, 238, 242)     # Blueprint highlighted block

    def set_slide_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_box(slide, left, top, width, height, text="", bg=BOX_BG, border=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
        
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Consolas"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
        return shape

    def add_header(slide, title, page_info="NEXUS GAME TRACKER - BLUEPRINT SPECIFICATION"):
        # Header bar
        add_box(slide, Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8), "", bg=RGBColor(240, 242, 245), border=BORDER_COLOR)
        
        # System Title / Page
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(7.0), Inches(0.7))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.name = "Consolas"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        # Subtitle
        tb2 = slide.shapes.add_textbox(Inches(7.8), Inches(0.5), Inches(4.7), Inches(0.6))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = page_info
        p2.alignment = PP_ALIGN.RIGHT
        p2.font.name = "Consolas"
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_MUTED

    def add_footer(slide, current_page, total_pages=15):
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.133), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"[ ARCHITECTURAL WIREFRAME BLUEPRINT ] | PAGE {current_page} OF {total_pages} | ISO-9241 UX COMPLIANT"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 1: Title & System Architecture
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1)
    
    # Outer frame
    add_box(s1, Inches(0.6), Inches(0.6), Inches(12.133), Inches(6.3), bg=RGBColor(255, 255, 255), border=BORDER_COLOR)
    
    tb = s1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(2.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "NEXUS GAME TRACKER"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Consolas"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    p2 = tf.add_paragraph()
    p2.text = "COMPLETE SYSTEM WIREFRAME BLUEPRINT (14 PAGES)"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Consolas"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED
    
    # Grid of Specs
    specs = [
        ("CLIENT PORTAL (7 SCREENS)", "index.html, catalog.html, client_search.html, client_logbook.html, client_profile.html, client_contact.html, login.html"),
        ("ADMIN TELEMETRY (7 SCREENS)", "admin.html, admin_games.html, admin_users.html, admin_analytics.html, admin_feedback.html, admin_logs.html, admin_settings.html"),
        ("DATA ARCHITECTURE", "Vanilla HTML5 + Modern CSS Design Tokens + LocalStorage Persistence + XML Seed"),
        ("DESIGN SYSTEM", "Monochromatic Architectural Grid, Strict Responsive Layouts, WCAG 2.2 AAA Contrast")
    ]
    for i, (head, desc) in enumerate(specs):
        x = Inches(1.5) if i % 2 == 0 else Inches(6.8)
        y = Inches(3.6) if i < 2 else Inches(4.9)
        b = add_box(s1, x, y, Inches(5.0), Inches(1.0), bg=ACCENT_BOX)
        tf_b = b.text_frame
        p_h = tf_b.paragraphs[0]
        p_h.text = head
        p_h.font.bold = True
        p_h.font.size = Pt(11)
        p_h.font.name = "Consolas"
        p_d = tf_b.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9)
        p_d.font.name = "Consolas"
        p_d.font.color.rgb = TEXT_MUTED

    add_footer(s1, 1)

    # ----------------------------------------------------
    # SLIDE 2: Index / Landing Page
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2)
    add_header(s2, "Layout 1: Landing Portal (index.html)")
    
    # Global Nav
    add_box(s2, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.6), "[LOGO] NEXUS               [CATALOG]   [SEARCH]   [VAULT]   [PROFILE]   [CONTACT]           [LOGIN / SIGN IN]", bg=ACCENT_BOX)
    
    # Hero Section
    add_box(s2, Inches(0.6), Inches(2.2), Inches(12.133), Inches(3.2), 
            "[ HERO BANNER WITH AMBIENT PARTICLES ]\n\nTRACK. DISCOVER. DOMINATE.\nYour Unified Gaming Telemetry & Personal Vault\n\n[ + ENTER VAULT CTA ]      [ EXPLORE CATALOG CTA ]")
    
    # 3 Feature Pillars
    add_box(s2, Inches(0.6), Inches(5.6), Inches(3.8), Inches(1.0), "[FEATURE 1]\nReal-Time RAWG Discovery")
    add_box(s2, Inches(4.766), Inches(5.6), Inches(3.8), Inches(1.0), "[FEATURE 2]\nInteractive Kanban Vault")
    add_box(s2, Inches(8.933), Inches(5.6), Inches(3.8), Inches(1.0), "[FEATURE 3]\nHardware Telemetry & Stats")
    add_footer(s2, 2)

    # ----------------------------------------------------
    # SLIDE 3: Catalog Page
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3)
    add_header(s3, "Layout 2: Discovery Hub (catalog.html)")
    add_box(s3, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.6), "[NAVBAR: Catalog Active]                                                     [Theme Toggle: Mocha / Latte]", bg=ACCENT_BOX)
    
    # Search & Filter bar
    add_box(s3, Inches(0.6), Inches(2.1), Inches(8.5), Inches(0.6), "[ SEARCH INPUT: 'Search 800,000+ Titles via RAWG Engine...' ]")
    add_box(s3, Inches(9.3), Inches(2.1), Inches(3.433), Inches(0.6), "[ FILTER: All / Action / RPG / Indie ]")
    
    # 4 Game Cards Grid
    for idx, name in enumerate(["Cyberpunk 2077", "Elden Ring", "Hades II", "Hollow Knight"]):
        x = Inches(0.6 + idx * 3.1)
        add_box(s3, x, Inches(2.9), Inches(2.833), Inches(3.7), 
                f"[ COVER ART IMAGE ]\n\n{name}\nGenre: Action / RPG\nRating: 4.8 / 5.0\n\n[ + ADD TO VAULT ]\n[ VIEW DETAILS ]")
    add_footer(s3, 3)

    # ----------------------------------------------------
    # SLIDE 4: Search Engine Page
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4)
    add_header(s4, "Layout 3: Telemetry Search (client_search.html)")
    add_box(s4, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.6), "[NAVBAR: Search Active]                                                              [Operative Status: Online]", bg=ACCENT_BOX)
    
    # Search Input Header
    add_box(s4, Inches(0.6), Inches(2.2), Inches(12.133), Inches(1.0), 
            "[ BIG SEARCH BAR: Enter Keywords, Tags, or Developers... ]\n[ ] Action  [ ] RPG  [ ] Strategy  [ ] Adventure  [ ] Platformer   |   Sort: [ Popularity ▼ ]")
    
    # Results Grid (2 Columns)
    add_box(s4, Inches(0.6), Inches(3.4), Inches(5.9), Inches(1.5), "[THUMB] Elden Ring (2022)\nFromSoftware | RPG, Souls-like\nMetacritic: 96 | [ + Quick Add ]")
    add_box(s4, Inches(6.833), Inches(3.4), Inches(5.9), Inches(1.5), "[THUMB] Cyberpunk 2077 (2020)\nCD Projekt Red | RPG, Open World\nMetacritic: 86 | [ + Quick Add ]")
    add_box(s4, Inches(0.6), Inches(5.1), Inches(5.9), Inches(1.5), "[THUMB] Hollow Knight (2017)\nTeam Cherry | Metroidvania\nMetacritic: 90 | [ + Quick Add ]")
    add_box(s4, Inches(6.833), Inches(5.1), Inches(5.9), Inches(1.5), "[THUMB] Hades II (2024)\nSupergiant Games | Roguelike\nEarly Access | [ + Quick Add ]")
    add_footer(s4, 4)

    # ----------------------------------------------------
    # SLIDE 5: Logbook / Kanban Vault
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5)
    add_header(s5, "Layout 4: Personal Vault & Kanban (client_logbook.html)")
    add_box(s5, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.6), "[NAVBAR: Vault Active]               [Total Games: 6]  [Hours Logged: 184 hrs]               [ + LOG NEW GAME ]", bg=ACCENT_BOX)
    
    # 3 Kanban Columns
    lanes = [
        ("IN PROGRESS / PLAYING (2)", ["[CARD] Cyberpunk 2077\n62 hrs | Rating: 5/5", "[CARD] Elden Ring\n122 hrs | Rating: 5/5"]),
        ("BACKLOG / QUEUED (3)", ["[CARD] Hades II\n0 hrs | Rating: N/A", "[CARD] Hollow Knight\n0 hrs | Rating: N/A", "[CARD] GTA VI\nPending Release"]),
        ("COMPLETED / ARCHIVED (1)", ["[CARD] Fable\n45 hrs | Rating: 4/5\n[ Status: Complete ]"])
    ]
    for idx, (lane_title, cards) in enumerate(lanes):
        x = Inches(0.6 + idx * 4.15)
        # Column Header
        add_box(s5, x, Inches(2.2), Inches(3.833), Inches(0.5), lane_title, bg=ACCENT_BOX)
        # Column Body Container
        add_box(s5, x, Inches(2.8), Inches(3.833), Inches(3.8), "", bg=BOX_BG)
        # Inner Cards
        for c_idx, card_text in enumerate(cards):
            add_box(s5, x + Inches(0.2), Inches(3.0 + c_idx * 1.15), Inches(3.433), Inches(1.0), card_text, bg=RGBColor(255, 255, 255))
    add_footer(s5, 5)

    # ----------------------------------------------------
    # SLIDE 6: Profile Page
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6)
    add_header(s6, "Layout 5: Operative Dossier (client_profile.html)")
    add_box(s6, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.6), "[NAVBAR: Profile Active]                                                          [Security Clearance: Operative Level 3]", bg=ACCENT_BOX)
    
    # Left: Identity Card
    add_box(s6, Inches(0.6), Inches(2.2), Inches(3.8), Inches(4.4), 
            "[ AVATAR BADGE ]\n\nCallsign: NEX_OPERATIVE\nStatus: Active Vanguard\nRank: Grandmaster (Tier V)\n\nSystem ID: #01CE1306\nJoined: August 2026\n\n[ EDIT PROFILE ]\n[ EXPORT JSON BACKUP ]")
    
    # Right Top: 3 Stat KPI Cards
    add_box(s6, Inches(4.6), Inches(2.2), Inches(2.5), Inches(1.2), "TOTAL TITLES\n6 Games")
    add_box(s6, Inches(7.3), Inches(2.2), Inches(2.5), Inches(1.2), "TIME IN COMBAT\n184 Hours")
    add_box(s6, Inches(10.0), Inches(2.2), Inches(2.733), Inches(1.2), "AVG SCORE\n4.8 / 5.0")
    
    # Right Bottom: Genre Telemetry & Recent Activity
    add_box(s6, Inches(4.6), Inches(3.6), Inches(4.0), Inches(3.0), "[ GENRE DISTRIBUTION CHART ]\n\n■ Action / RPG: 50%\n■ Metroidvania: 25%\n■ Open World: 25%")
    add_box(s6, Inches(8.8), Inches(3.6), Inches(3.933), Inches(3.0), "[ RECENT TELEMETRY LOG ]\n\n• Logged 4 hrs on Elden Ring\n• Moved Hades II to Backlog\n• Seeded Library from XML\n• Synced LocalStorage State")
    add_footer(s6, 6)

    # ----------------------------------------------------
    # SLIDE 7: Contact Page
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7)
    add_header(s7, "Layout 6: Support & Feedback Uplink (client_contact.html)")
    add_box(s7, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.6), "[NAVBAR: Contact Active]                                                          [Support Channel: Encrypted TLS 1.3]", bg=ACCENT_BOX)
    
    # Left: Dispatch Guidelines
    add_box(s7, Inches(0.6), Inches(2.2), Inches(4.5), Inches(4.4), 
            "[ TRANSMISSION DISPATCH ]\n\nHave feedback, bug reports, or feature requests?\n\nDirect Inquiries:\n• Email: support@nexus.local\n• Matrix: @nexus:local.net\n• Response Window: < 24 Hours\n\nAll transmissions are logged directly to the Admin Support Queue.")
    
    # Right: Form Container
    add_box(s7, Inches(5.3), Inches(2.2), Inches(7.433), Inches(4.4), 
            "[ TRANSMISSION FORM ]\n\nOperative Callsign Input: [ e.g. Spectre-01 ]\n\nEmail Uplink: [ operative@nexus.local ]\n\nCategory: [ Bug Report ▼ ]\n\nMessage Payload:\n[ Enter detailed feedback or telemetry logs... ]\n\n[ >>> TRANSMIT MESSAGE PAYLOAD <<< ]")
    add_footer(s7, 7)

    # ----------------------------------------------------
    # SLIDE 8: Authentication Page
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8)
    add_header(s8, "Layout 7: Terminal Authentication (login.html)")
    add_box(s8, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.6), "[NAVBAR: Auth Terminal]                                                           [Protocol: SHA-256 Client-Side Token]", bg=ACCENT_BOX)
    
    # Dual Auth Containers
    add_box(s8, Inches(1.5), Inches(2.3), Inches(4.8), Inches(4.2), 
            "[ OPERATIVE SIGN IN ]\n\nCallsign / Username:\n[ input_field ]\n\nPassphrase Key:\n[ password_field ]\n\n[ ✓ ] Remember Operative State\n\n[ >>> AUTHENTICATE & ENTER <<< ]\n\nDemo: operative / pass123")
    
    add_box(s8, Inches(7.0), Inches(2.3), Inches(4.8), Inches(4.2), 
            "[ NEW OPERATIVE REGISTRATION ]\n\nDesired Callsign:\n[ input_field ]\n\nEmail Uplink:\n[ email_field ]\n\nMaster Passphrase:\n[ password_field ]\n\n[ >>> INITIALIZE NEW PROFILE <<< ]\n\nInstant LocalStorage Account Creation")
    add_footer(s8, 8)

    # ----------------------------------------------------
    # SLIDE 9: Admin Dashboard (Mission Control)
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9)
    add_header(s9, "Layout 8: Admin Mission Control (admin.html)", "ADMINISTRATOR PRIVILEGES ENABLED")
    
    # Sidebar + Content layout
    add_box(s9, Inches(0.6), Inches(1.4), Inches(2.6), Inches(5.2), 
            "[ ADMIN SIDEBAR ]\n\n■ Dashboard (Active)\n■ Games Registry\n■ Operatives DB\n■ Telemetry Analytics\n■ Support Inbox\n■ Security Audit Logs\n■ System Settings\n\n[ LOGOUT ]", bg=ACCENT_BOX)
    
    # 4 Admin KPI Cards
    add_box(s9, Inches(3.4), Inches(1.4), Inches(2.15), Inches(1.1), "TOTAL GAMES\n6 Loaded")
    add_box(s9, Inches(5.7), Inches(1.4), Inches(2.15), Inches(1.1), "OPERATIVES\n3 Registered")
    add_box(s9, Inches(8.0), Inches(1.4), Inches(2.15), Inches(1.1), "INBOX MSGS\n2 Pending")
    add_box(s9, Inches(10.3), Inches(1.4), Inches(2.433), Inches(1.1), "NODE HEALTH\n100% Online")
    
    # Big Chart Box & Recent Feeds
    add_box(s9, Inches(3.4), Inches(2.7), Inches(5.5), Inches(3.9), "[ SYSTEM ANALYTICS & TRAFFIC TELEMETRY ]\n\n• Weekly Session Hours: 184 hrs\n• Storage Used: 48 KB / 5 MB (LocalStorage)\n• XML Pipeline Sync: Operational\n• Theme Switch Frequency: Mocha 65% / Latte 35%")
    add_box(s9, Inches(9.1), Inches(2.7), Inches(3.633), Inches(3.9), "[ LIVE AUDIT STREAM ]\n\n[INFO] Admin logged in\n[INFO] Games seeded from XML\n[INFO] Modal CRUD form initialized\n[WARN] API Rate limit: 85% remaining\n[INFO] Contact payload stored")
    add_footer(s9, 9)

    # ----------------------------------------------------
    # SLIDE 10: Admin Games Registry
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10)
    add_header(s10, "Layout 9: Admin Game Registry & CRUD (admin_games.html)", "ADMINISTRATOR PRIVILEGES ENABLED")
    
    add_box(s10, Inches(0.6), Inches(1.4), Inches(2.6), Inches(5.2), "[ SIDEBAR ]\n\n■ Dashboard\n■ Games (Active)\n■ Operatives\n■ Analytics\n■ Support Inbox\n■ Audit Logs\n■ Settings", bg=ACCENT_BOX)
    
    # Header bar with Actions
    add_box(s10, Inches(3.4), Inches(1.4), Inches(5.8), Inches(0.6), "[ SEARCH REGISTRY: Filter by title, genre, year... ]")
    add_box(s10, Inches(9.4), Inches(1.4), Inches(3.333), Inches(0.6), "[ + ADD NEW GAME MODAL ]", bg=ACCENT_BOX)
    
    # Data Table
    table_header = "COVER    TITLE              GENRE          STATUS       HOURS    SCORE    ACTIONS"
    r1 = "[IMG]    Elden Ring         RPG            Playing      122 hrs  5/5      [Edit] [Delete]"
    r2 = "[IMG]    Cyberpunk 2077     Action/RPG     Playing      62 hrs   5/5      [Edit] [Delete]"
    r3 = "[IMG]    Hollow Knight      Metroidvania   Backlog      0 hrs    5/5      [Edit] [Delete]"
    r4 = "[IMG]    Hades II           Roguelike      Backlog      0 hrs    N/A      [Edit] [Delete]"
    
    add_box(s10, Inches(3.4), Inches(2.2), Inches(9.333), Inches(4.4), 
            f"{table_header}\n" + "-"*75 + f"\n{r1}\n{r2}\n{r3}\n{r4}\n\n[ MODAL OVERLAYS: #gameFormModal (Add/Edit) | #deleteConfirmModal (Delete) ]")
    add_footer(s10, 10)

    # ----------------------------------------------------
    # SLIDE 11: Admin Operatives Management
    # ----------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s11)
    add_header(s11, "Layout 10: Admin Operatives Roster (admin_users.html)", "ADMINISTRATOR PRIVILEGES ENABLED")
    add_box(s11, Inches(0.6), Inches(1.4), Inches(2.6), Inches(5.2), "[ SIDEBAR: Users Active ]", bg=ACCENT_BOX)
    
    add_box(s11, Inches(3.4), Inches(1.4), Inches(9.333), Inches(0.6), "[ SEARCH OPERATIVES ]                                                    [ + INVITE OPERATIVE ]")
    
    u_head = "UID      CALLSIGN          ROLE             STATUS       GAMES LOGGED    LAST ACCESS    ACTIONS"
    u1 =     "#001     NexBoss           Root Admin       Active       12 Titles       Just Now       [Edit] [Audit]"
    u2 =     "#002     ItzJazzu          Lead Vanguard    Active       6 Titles        Today          [Edit] [Audit]"
    u3 =     "#003     CyberSpectre      Operative        Standby      4 Titles        Yesterday      [Edit] [Demote]"
    
    add_box(s11, Inches(3.4), Inches(2.2), Inches(9.333), Inches(4.4), 
            f"{u_head}\n" + "-"*82 + f"\n{u1}\n{u2}\n{u3}\n\n[ ROLE HIERARCHY: Root Administrator > Vanguard Operative > Guest Recon ]")
    add_footer(s11, 11)

    # ----------------------------------------------------
    # SLIDE 12: Admin Analytics
    # ----------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s12)
    add_header(s12, "Layout 11: System Telemetry & Analytics (admin_analytics.html)", "ADMINISTRATOR PRIVILEGES ENABLED")
    add_box(s12, Inches(0.6), Inches(1.4), Inches(2.6), Inches(5.2), "[ SIDEBAR: Analytics Active ]", bg=ACCENT_BOX)
    
    add_box(s12, Inches(3.4), Inches(1.4), Inches(4.5), Inches(2.5), "[ PLAYTIME VELOCITY (HOURS/WEEK) ]\n\n■ Elden Ring: 66%\n■ Cyberpunk 2077: 34%\n\nTotal System Playtime: 184 Hours")
    add_box(s12, Inches(8.1), Inches(1.4), Inches(4.633), Inches(2.5), "[ PLATFORM DISTRIBUTION ]\n\n■ PC / Steam: 65%\n■ PlayStation 5: 20%\n■ Xbox Series X: 15%")
    add_box(s12, Inches(3.4), Inches(4.1), Inches(9.333), Inches(2.5), "[ RAWG API TRAFFIC & LOCALSTORAGE CAPACITY ]\n\n• Cache Hit Rate: 94.2%\n• Average Response Time: 42ms\n• Storage Footprint: 48.2 KB (1.0% quota)\n• Active Theme Distribution: Catppuccin Mocha (62%) / Latte (38%)")
    add_footer(s12, 12)

    # ----------------------------------------------------
    # SLIDE 13: Admin Feedback / Support Inbox
    # ----------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s13)
    add_header(s13, "Layout 12: Support Inbox & Transmissions (admin_feedback.html)", "ADMINISTRATOR PRIVILEGES ENABLED")
    add_box(s13, Inches(0.6), Inches(1.4), Inches(2.6), Inches(5.2), "[ SIDEBAR: Inbox Active ]", bg=ACCENT_BOX)
    
    add_box(s13, Inches(3.4), Inches(1.4), Inches(9.333), Inches(0.6), "[ INBOX FILTER: All Transmissions (2) / Unread (1) / Resolved (1) ]")
    
    m_head = "ID       FROM CALLSIGN         EMAIL                    SUBJECT / PAYLOAD               TIME       ACTION"
    m1 =     "#TR-01   Vanguard-9            vanguard9@nexus.local    'Steam API sync request'        10m ago    [Reply] [Resolve]"
    m2 =     "#TR-02   ShadowRecon           recon@nexus.local        'Great Catppuccin theme!'       2h ago     [Archived]"
    
    add_box(s13, Inches(3.4), Inches(2.2), Inches(9.333), Inches(4.4), 
            f"{m_head}\n" + "-"*82 + f"\n{m1}\n{m2}\n\n[ MODAL OVERLAY: #viewTransmissionModal | Reply Dispatch Engine ]")
    add_footer(s13, 13)

    # ----------------------------------------------------
    # SLIDE 14: Admin Security Audit Logs
    # ----------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s14)
    add_header(s14, "Layout 13: Security Audit Ledger (admin_logs.html)", "ADMINISTRATOR PRIVILEGES ENABLED")
    add_box(s14, Inches(0.6), Inches(1.4), Inches(2.6), Inches(5.2), "[ SIDEBAR: Audit Logs Active ]", bg=ACCENT_BOX)
    
    l_head = "TIMESTAMP            SEVERITY    OPERATIVE       EVENT DESCRIPTION"
    l1 =     "2026-08-30 09:20:46  [INFO]      System          XML database pre-seeded into LocalStorage (6 titles)"
    l2 =     "2026-08-30 09:20:50  [INFO]      NexBoss         Switched theme palette to Catppuccin Mocha"
    l3 =     "2026-08-30 09:21:07  [INFO]      Spectre         Submitted support transmission #TR-01"
    l4 =     "2026-08-30 09:21:15  [SUCCESS]   Git Engine      Committed build 2c7da8f to production Netlify pipeline"
    
    add_box(s14, Inches(3.4), Inches(1.4), Inches(9.333), Inches(5.2), 
            f"[ REAL-TIME SECURITY LEDGER ]\n\n{l_head}\n" + "-"*75 + f"\n{l1}\n{l2}\n{l3}\n{l4}\n\n[ EXPORT AUDIT LEDGER (CSV / JSON) ]")
    add_footer(s14, 14)

    # ----------------------------------------------------
    # SLIDE 15: Admin System Settings
    # ----------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s15)
    add_header(s15, "Layout 14: System Settings (admin_settings.html)", "ADMINISTRATOR PRIVILEGES ENABLED")
    add_box(s15, Inches(0.6), Inches(1.4), Inches(2.6), Inches(5.2), "[ SIDEBAR: Settings Active ]", bg=ACCENT_BOX)
    
    add_box(s15, Inches(3.4), Inches(1.4), Inches(9.333), Inches(5.2), 
            "[ GLOBAL CONFIGURATION PARAMETERS ]\n\n"
            "• Platform Title:       [ Nexus Game Tracker (01CE1306) ]\n"
            "• RAWG API Key:         [ ******************************** ]\n"
            "• Default Theme:        (●) Catppuccin Mocha     (○) Catppuccin Latte\n"
            "• Auto-Seed XML:        [✓] Enabled on first user session\n"
            "• Glassmorphism FX:     [✓] Enabled (Backdrop Filter: blur(12px))\n"
            "• Security Level:       [ High (Strict LocalStorage Isolation) ▼ ]\n\n"
            "[ >>> SAVE ALL CONFIGURATION CHANGES <<< ]      [ RESET TO FACTORY XML SEED ]")
    add_footer(s15, 15)

    output_path = ".dev_workspace/Nexus_Wireframes_BlackWhite.pptx"
    prs.save(output_path)
    print(f"Successfully generated 15-slide B&W Blueprint wireframe at: {output_path}")

if __name__ == "__main__":
    create_bw_blueprint_pptx()
